# -*- coding: utf-8 -*-
"""
일반 Why→What→How 문서( 07-적용_CDSEM_실행기획/ 의 11단계 등 )  →  PPTX 변환기.

`md2pptx.py`(장표 전용)와 달리, 슬라이드 경계를 `## 장표 N` 이 아니라
**문서의 헤더 구조**로 잡는다:

- 입력: 마크다운 파일 여러 개(순서대로) + 출력 경로
- 각 파일:
    · frontmatter(--- … ---) 제거
    · `# H1` + 바로 뒤 `>` 인용  →  섹션 구분(디바이더) 슬라이드 1장
    · `## H2`  →  슬라이드 1장 (하위 `### H3` 가 있으면 H3마다 추가 분할)
    · `## 이전/다음 단계` 류 네비게이션 섹션은 건너뜀
- 표/코드펜스/블록쿼트/불릿 렌더링은 md2pptx.py 와 동일 스타일.

사용:  python tools/md2pptx_doc.py <out.pptx> <in1.md> <in2.md> ...
       (입력 생략 시 07 폴더 01~11 자동 수집)
한글 폰트: 맑은 고딕.
"""
import glob
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

KR_FONT = "맑은 고딕"
MONO_FONT = "Consolas"
ACCENT = RGBColor(0x1F, 0x49, 0x7D)      # 진청 (헤더/타이틀)
ACCENT_LT = RGBColor(0xDC, 0xE6, 0xF1)   # 연청 (표 헤더 배경)
GREY = RGBColor(0x59, 0x59, 0x59)
DARK = RGBColor(0x26, 0x26, 0x26)
CODE_BG = RGBColor(0xF2, 0xF4, 0xF7)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_W = Emu(SLIDE_W - 2 * MARGIN)
CONTENT_TOP = Inches(1.35)
CONTENT_BOTTOM = Inches(7.2)

# 건너뛸 섹션(순수 네비게이션·메타)
SKIP_H2 = ("이전/다음 단계", "이전 / 다음", "참고", "참고 자료")


# ---------- 인라인 포맷 ----------
LINK_RE = re.compile(r"\[([^\]]+)\]\([^)]+\)")


def strip_inline(text):
    text = LINK_RE.sub(r"\1", text)
    text = text.replace("`", "")
    text = text.replace("<br>", "\n")
    return text.strip()


def add_runs(paragraph, text, size, color=DARK, base_bold=False, font=KR_FONT):
    parts = re.split(r"(\*\*)", text)
    bold = base_bold
    for p in parts:
        if p == "**":
            bold = not bold
            continue
        if p == "":
            continue
        p = p.replace("*", "")
        run = paragraph.add_run()
        run.text = p
        run.font.size = size
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color


def set_cell_text(cell, text, size, color=DARK, bold=False, header=False):
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    lines = strip_inline(text).split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        add_runs(para, line, size, color=color, base_bold=bold or header)
    if header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_LT


# ---------- 마크다운 → 블록 파서 ----------
def parse_blocks(lines):
    """라인들 → 타입 블록 리스트. (heading 은 level 보존)"""
    blocks = []
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        s = line.strip()
        if s == "":
            i += 1
            continue
        if re.match(r"^-{3,}$", s):
            i += 1
            continue
        # 코드펜스
        if s.startswith("```"):
            code = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i].rstrip("\n"))
                i += 1
            i += 1
            blocks.append(("code", "\n".join(code)))
            continue
        # 표
        if s.startswith("|"):
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                rows.append(lines[i].strip())
                i += 1
            table = []
            for r in rows:
                if re.match(r"^\|[\s:\-\|]+\|$", r):
                    continue
                cells = [c.strip() for c in r.strip("|").split("|")]
                table.append(cells)
            if table:
                blocks.append(("table", table))
            continue
        # 블록쿼트
        if s.startswith(">"):
            quote = []
            while i < n and lines[i].strip().startswith(">"):
                quote.append(lines[i].strip().lstrip(">").strip())
                i += 1
            text = "\n".join(quote).strip()
            blocks.append(("callout", text))
            continue
        # 헤더
        if s.startswith("#"):
            lvl = len(s) - len(s.lstrip("#"))
            blocks.append(("heading", s.lstrip("#").strip(), lvl))
            i += 1
            continue
        # 불릿 / 번호
        m = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)$", line)
        if m:
            indent = len(m.group(1))
            level = 1 if indent >= 2 else 0
            ordered = bool(re.match(r"\d+\.", m.group(2)))
            prefix = (m.group(2) + " ") if ordered else ""
            blocks.append(("bullet", prefix + m.group(3), level))
            i += 1
            continue
        blocks.append(("para", s))
        i += 1
    return blocks


# ---------- 높이 추정 ----------
def est_text_height(text, size_pt, width_in, font_factor=1.9):
    chars_per_line = max(8, int(width_in * 72 / (size_pt * 0.62)))
    total = 0
    for ln in text.split("\n"):
        total += max(1, -(-len(ln) // chars_per_line))
    line_h = size_pt * font_factor
    return Pt(total * line_h)


# ---------- 슬라이드 렌더 ----------
def add_title(slide, title, idx=""):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    box = slide.shapes.add_textbox(MARGIN, Inches(0.30), CONTENT_W - Inches(1.0), Inches(1.0))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_runs(p, title, Pt(22), color=ACCENT, base_bold=True)
    if idx:
        num = slide.shapes.add_textbox(SLIDE_W - Inches(1.3), Inches(0.35), Inches(1.0), Inches(0.5))
        np = num.text_frame.paragraphs[0]; np.alignment = PP_ALIGN.RIGHT
        r = np.add_run(); r.text = str(idx); r.font.size = Pt(13); r.font.color.rgb = GREY; r.font.name = KR_FONT


def render_table(slide, table_data, top):
    rows = len(table_data)
    cols = max(len(r) for r in table_data)
    for r in table_data:
        while len(r) < cols:
            r.append("")
    body_size = 10
    head_size = 10.5
    if cols == 2:
        widths = [CONTENT_W * 0.28, CONTENT_W * 0.72]
    else:
        widths = [int(CONTENT_W / cols)] * cols
    row_heights = []
    for ri, row in enumerate(table_data):
        max_lines = 1
        for ci, cell in enumerate(row):
            w_in = Emu(int(widths[ci])).inches
            txt = strip_inline(cell)
            cpl = max(6, int(w_in * 72 / ((head_size if ri == 0 else body_size) * 0.62)))
            lines = 0
            for ln in txt.split("\n"):
                lines += max(1, -(-len(ln) // cpl))
            max_lines = max(max_lines, lines)
        row_heights.append(Pt(max_lines * (body_size * 1.5) + 5))
    total_h = sum(int(h) for h in row_heights)
    gfx = slide.shapes.add_table(rows, cols, MARGIN, top, CONTENT_W, total_h)
    tbl = gfx.table
    tbl.first_row = True
    for ci in range(cols):
        tbl.columns[ci].width = int(widths[ci])
    for ri in range(rows):
        tbl.rows[ri].height = int(row_heights[ri])
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            is_head = (ri == 0)
            set_cell_text(cell, table_data[ri][ci],
                          Pt(head_size if is_head else body_size),
                          color=ACCENT if is_head else DARK,
                          header=is_head)
            if not is_head:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return top + Emu(total_h) + Inches(0.12)


def render_code(slide, code, top):
    size = 8.5
    h = est_text_height(code, size, CONTENT_W.inches, font_factor=1.4)
    box = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, h + Inches(0.16))
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0xCF, 0xD6, 0xDF); box.line.width = Pt(0.5)
    tf = box.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = MONO_FONT; r.font.size = Pt(size); r.font.color.rgb = DARK
    return top + Emu(int(h)) + Inches(0.28)


def render_text(slide, kind, text, top, level=0):
    if kind == "heading":
        size, color, bold = 14, ACCENT, True
    elif kind == "callout":
        size, color, bold = 11, RGBColor(0x33, 0x44, 0x55), False
    elif kind == "bullet":
        size, color, bold = 11.5, DARK, False
    else:
        size, color, bold = 11.5, DARK, False
    h = est_text_height(strip_inline(text), size, CONTENT_W.inches - level * 0.3)
    box = slide.shapes.add_textbox(MARGIN + Inches(0.05 + level * 0.3), top,
                                   CONTENT_W - Inches(level * 0.3), h + Inches(0.06))
    if kind == "callout":
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFB)
        box.line.color.rgb = ACCENT_LT; box.line.width = Pt(1)
        box.text_frame.margin_left = Inches(0.12)
    tf = box.text_frame; tf.word_wrap = True
    for i, ln in enumerate(strip_inline(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = ""
        if kind == "bullet":
            prefix = ("    • " if level else "• ")
            if re.match(r"^\d+\.", ln):
                prefix = ("    " if level else "")
        if prefix:
            r = p.add_run(); r.text = prefix; r.font.size = Pt(size); r.font.name = KR_FONT; r.font.color.rgb = color
        add_runs(p, ln, Pt(size), color=color, base_bold=bold)
    return top + Emu(int(h)) + Inches(0.10)


def render_blocks(slide, blocks):
    y = CONTENT_TOP
    for b in blocks:
        if b[0] == "table":
            y = render_table(slide, b[1], y)
        elif b[0] == "code":
            y = render_code(slide, b[1], y)
        elif b[0] == "heading":
            y = render_text(slide, "heading", b[1], y)
        elif b[0] == "bullet":
            y = render_text(slide, "bullet", b[1], y, level=b[2])
        elif b[0] == "callout":
            y = render_text(slide, "callout", b[1], y)
        else:
            y = render_text(slide, "para", b[1], y)


def add_divider(prs, step, title, subtitle):
    """파일 시작 — 큰 단계 제목 + 한 줄 목표."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = ACCENT; bg.line.fill.background()
    box = slide.shapes.add_textbox(MARGIN, Inches(2.4), CONTENT_W, Inches(1.4))
    tf = box.text_frame; tf.word_wrap = True
    if step:
        p0 = tf.paragraphs[0]
        r = p0.add_run(); r.text = f"STEP {step}"
        r.font.size = Pt(20); r.font.bold = True; r.font.name = KR_FONT
        r.font.color.rgb = RGBColor(0xBF, 0xD3, 0xE6)
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    add_runs(p, title, Pt(30), color=RGBColor(0xFF, 0xFF, 0xFF), base_bold=True)
    if subtitle:
        sb = slide.shapes.add_textbox(MARGIN, Inches(4.0), CONTENT_W, Inches(2.2))
        stf = sb.text_frame; stf.word_wrap = True
        for i, ln in enumerate(strip_inline(subtitle).split("\n")):
            sp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
            add_runs(sp, ln, Pt(14), color=RGBColor(0xE6, 0xEE, 0xF5))


def split_frontmatter(md):
    if md.startswith("---"):
        m = re.match(r"^---\n.*?\n---\n", md, re.DOTALL)
        if m:
            return md[m.end():]
    return md


def process_file(prs, path):
    with open(path, encoding="utf-8") as f:
        md = split_frontmatter(f.read())
    lines = md.split("\n")

    # H1 + 직후 인용 추출
    h1 = os.path.basename(path)
    intro = []
    body_start = 0
    for i, ln in enumerate(lines):
        if ln.startswith("# "):
            h1 = ln[2:].strip()
            j = i + 1
            while j < len(lines) and (lines[j].strip() == "" or lines[j].strip().startswith(">")):
                if lines[j].strip().startswith(">"):
                    intro.append(lines[j].strip().lstrip(">").strip())
                j += 1
            body_start = j
            break

    # 파일명에서 단계 번호
    base = os.path.basename(path)
    mstep = re.match(r"^(\d+)", base)
    step = mstep.group(1) if mstep else ""
    # 디바이더 제목: H1 에서 "NN · " 접두 제거
    div_title = re.sub(r"^\d+\s*[·.]\s*", "", h1)

    add_divider(prs, step, div_title, "\n".join(intro))

    # 본문을 ## 단위로 슬라이드 분할 (### 가 있으면 추가 분할)
    body = lines[body_start:]
    # 섹션 단위로 자르기
    sections = []  # (h2title, [lines])
    cur_title = None
    cur = []
    for ln in body:
        if ln.startswith("## "):
            if cur_title is not None or cur:
                sections.append((cur_title, cur))
            cur_title = ln[3:].strip()
            cur = []
        else:
            cur.append(ln)
    if cur_title is not None or cur:
        sections.append((cur_title, cur))

    for h2title, seclines in sections:
        if h2title is None:
            continue
        if any(h2title.startswith(s) for s in SKIP_H2):
            continue
        # ### 로 다시 분할
        sub = []  # (subtitle_or_None, [lines])
        sub_title = None
        sub_lines = []
        for ln in seclines:
            if ln.startswith("### "):
                if sub_title is not None or sub_lines:
                    sub.append((sub_title, sub_lines))
                sub_title = ln[4:].strip()
                sub_lines = []
            else:
                sub_lines.append(ln)
        if sub_title is not None or sub_lines:
            sub.append((sub_title, sub_lines))

        for sub_title, sub_lines in sub:
            blocks = parse_blocks(sub_lines)
            if not blocks:
                continue
            if sub_title:
                title = f"{step} · {h2title} — {sub_title}"
            else:
                title = f"{step} · {h2title}"
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_title(slide, title, idx=step)
            render_blocks(slide, blocks)


def main():
    if len(sys.argv) < 2:
        print("usage: python tools/md2pptx_doc.py <out.pptx> [in1.md in2.md ...]")
        sys.exit(1)
    out = sys.argv[1]
    inputs = sys.argv[2:]
    if not inputs:
        # 기본: 07 폴더 01~11
        base = os.path.join(os.path.dirname(out) or ".", "07-적용_CDSEM_실행기획")
        inputs = sorted(glob.glob(os.path.join(base, "[01][0-9]-*.md")))
        inputs = [p for p in inputs if not os.path.basename(p).startswith("00")]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for path in inputs:
        process_file(prs, path)

    prs.save(out)
    print(f"OK  files={len(inputs)}  slides={len(prs.slides._sldIdLst)}  ->  {out}")


if __name__ == "__main__":
    main()
