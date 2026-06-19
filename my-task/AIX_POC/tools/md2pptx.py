# -*- coding: utf-8 -*-
"""
06-적용_CDSEM_DesignCamp장표.md  →  PPTX 변환기.

장표 마크다운을 편집 가능한 네이티브 PPTX로 변환한다.
- `## 장표 N — Title`  = 슬라이드 1장
- 마크다운 표           = 네이티브 PPT 표
- 코드펜스(```)        = 고정폭(Consolas) 텍스트박스 (ASCII 다이어그램)
- `> 🎤 발표 스크립트`  = 발표자 노트(notes slide)
- 그 외 `>` 인용        = 본문 콜아웃
- 헤더(#/###)/불릿/문단 = 텍스트
한글 폰트: 맑은 고딕(Malgun Gothic).
"""
import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = sys.argv[1] if len(sys.argv) > 1 else "06-적용_CDSEM_DesignCamp장표.md"
OUT = sys.argv[2] if len(sys.argv) > 2 else "06-적용_CDSEM_DesignCamp장표.pptx"

KR_FONT = "맑은 고딕"
MONO_FONT = "Consolas"
ACCENT = RGBColor(0x1F, 0x49, 0x7D)      # 진청 (헤더/타이틀)
ACCENT_LT = RGBColor(0xDC, 0xE6, 0xF1)   # 연청 (표 헤더 배경)
GREY = RGBColor(0x59, 0x59, 0x59)
DARK = RGBColor(0x26, 0x26, 0x26)
CODE_BG = RGBColor(0xF2, 0xF4, 0xF7)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_W = Emu(SLIDE_W - 2 * MARGIN)
CONTENT_TOP = Inches(1.35)
CONTENT_BOTTOM = Inches(7.2)


# ---------- 인라인 포맷 ----------
LINK_RE = re.compile(r"\[([^\]]+)\]\([^)]+\)")

def strip_inline(text):
    """링크→앵커텍스트, 백틱 제거, <br>→\n, **는 보존(굵게용)."""
    text = LINK_RE.sub(r"\1", text)
    text = text.replace("`", "")
    text = text.replace("<br>", "\n")
    return text.strip()


def add_runs(paragraph, text, size, color=DARK, base_bold=False, font=KR_FONT):
    """**굵게** 토글을 런으로 분해. \n 은 호출 전에 분리."""
    parts = re.split(r"(\*\*)", text)
    bold = base_bold
    for p in parts:
        if p == "**":
            bold = not bold
            continue
        if p == "":
            continue
        # *이탤릭* 단순 제거 처리(강조만 굵게로 통일)
        p = p.replace("*", "")
        run = paragraph.add_run()
        run.text = p
        run.font.size = size
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color


def set_cell_text(cell, text, size, color=DARK, bold=False, header=False):
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    lines = strip_inline(text).split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        add_runs(para, line, size, color=color, base_bold=bold or header)
    if header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_LT


# ---------- 마크다운 → 블록 파서 ----------
def parse_blocks(lines):
    """한 슬라이드의 라인들 → 타입 블록 리스트 + notes 문자열."""
    blocks = []
    notes = []
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        s = line.strip()
        if s == "":
            i += 1
            continue
        # 수평 구분선(---) 은 슬라이드 경계라 무시
        if re.match(r"^-{3,}$", s):
            i += 1
            continue
        # 코드펜스
        if s.startswith("```"):
            code = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i].rstrip("\n"))
                i += 1
            i += 1  # 닫는 ```
            blocks.append(("code", "\n".join(code)))
            continue
        # 표
        if s.startswith("|"):
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                rows.append(lines[i].strip())
                i += 1
            table = []
            for r in rows:
                if re.match(r"^\|[\s:\-\|]+\|$", r):  # 구분행
                    continue
                cells = [c.strip() for c in r.strip("|").split("|")]
                table.append(cells)
            if table:
                blocks.append(("table", table))
            continue
        # 블록쿼트
        if s.startswith(">"):
            quote = []
            while i < n and lines[i].strip().startswith(">"):
                quote.append(lines[i].strip()[1:].strip())
                i += 1
            text = "\n".join(quote).strip()
            if "🎤" in text:
                # 발표 스크립트 → 노트
                t = text.replace("🎤", "").replace("**발표 스크립트**", "").replace("**", "").strip()
                notes.append(t)
            else:
                blocks.append(("callout", text))
            continue
        # 헤더
        if s.startswith("#"):
            lvl = len(s) - len(s.lstrip("#"))
            blocks.append(("heading", s.lstrip("#").strip(), lvl))
            i += 1
            continue
        # 불릿 / 번호
        m = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)$", line)
        if m:
            indent = len(m.group(1))
            level = 1 if indent >= 2 else 0
            ordered = bool(re.match(r"\d+\.", m.group(2)))
            prefix = (m.group(2) + " ") if ordered else ""
            blocks.append(("bullet", prefix + m.group(3), level))
            i += 1
            continue
        # 일반 문단
        blocks.append(("para", s))
        i += 1
    return blocks, "\n\n".join(notes)


# ---------- 높이 추정(스택 레이아웃) ----------
def est_text_height(text, size_pt, width_in, font_factor=1.9):
    """대략적 줄 수 × 줄높이. 한글 폭 보정."""
    chars_per_line = max(8, int(width_in * 72 / (size_pt * 0.62)))
    total = 0
    for ln in text.split("\n"):
        total += max(1, -(-len(ln) // chars_per_line))  # ceil
    line_h = size_pt * font_factor
    return Pt(total * line_h)


# ---------- 슬라이드 렌더 ----------
def add_title(slide, title, idx):
    # 상단 악센트 바
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    box = slide.shapes.add_textbox(MARGIN, Inches(0.32), CONTENT_W - Inches(1.0), Inches(0.95))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_runs(p, title, Pt(24), color=ACCENT, base_bold=True)
    # 우상단 페이지 번호
    num = slide.shapes.add_textbox(SLIDE_W - Inches(1.1), Inches(0.35), Inches(0.8), Inches(0.5))
    np = num.text_frame.paragraphs[0]; np.alignment = PP_ALIGN.RIGHT
    r = np.add_run(); r.text = str(idx); r.font.size = Pt(14); r.font.color.rgb = GREY; r.font.name = KR_FONT


def render_table(slide, table_data, top):
    rows = len(table_data)
    cols = max(len(r) for r in table_data)
    # 셀 보정
    for r in table_data:
        while len(r) < cols:
            r.append("")
    # 행 높이 추정
    body_size = 10.5
    head_size = 11
    # 열 너비: 2열 key-value면 1열을 좁게
    if cols == 2:
        widths = [CONTENT_W * 0.28, CONTENT_W * 0.72]
    else:
        widths = [int(CONTENT_W / cols)] * cols
    # 행 높이: 셀 내 최대 줄 수 기반
    row_heights = []
    for ri, row in enumerate(table_data):
        max_lines = 1
        for ci, cell in enumerate(row):
            w_in = Emu(int(widths[ci])).inches
            txt = strip_inline(cell)
            cpl = max(6, int(w_in * 72 / ((head_size if ri == 0 else body_size) * 0.62)))
            lines = 0
            for ln in txt.split("\n"):
                lines += max(1, -(-len(ln) // cpl))
            max_lines = max(max_lines, lines)
        row_heights.append(Pt(max_lines * (body_size * 1.55) + 6))
    total_h = sum(int(h) for h in row_heights)
    gfx = slide.shapes.add_table(rows, cols, MARGIN, top, CONTENT_W, total_h)
    tbl = gfx.table
    tbl.first_row = True
    for ci in range(cols):
        tbl.columns[ci].width = int(widths[ci])
    for ri in range(rows):
        tbl.rows[ri].height = int(row_heights[ri])
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            is_head = (ri == 0)
            set_cell_text(cell, table_data[ri][ci],
                          Pt(head_size if is_head else body_size),
                          color=ACCENT if is_head else DARK,
                          header=is_head)
            if not is_head:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return top + Emu(total_h) + Inches(0.12)


def render_code(slide, code, top):
    size = 9
    h = est_text_height(code, size, CONTENT_W.inches, font_factor=1.45)
    box = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, h + Inches(0.16))
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0xCF, 0xD6, 0xDF); box.line.width = Pt(0.5)
    tf = box.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = MONO_FONT; r.font.size = Pt(size); r.font.color.rgb = DARK
    return top + Emu(int(h)) + Inches(0.28)


def render_text(slide, kind, text, top, level=0):
    if kind == "heading":
        size, color, bold = 15, ACCENT, True
    elif kind == "callout":
        size, color, bold = 11.5, RGBColor(0x33, 0x44, 0x55), False
    elif kind == "bullet":
        size, color, bold = 12, DARK, False
    else:
        size, color, bold = 12, DARK, False
    h = est_text_height(strip_inline(text), size, CONTENT_W.inches - level * 0.3)
    box = slide.shapes.add_textbox(MARGIN + Inches(0.05 + level * 0.3), top,
                                   CONTENT_W - Inches(level * 0.3), h + Inches(0.06))
    if kind == "callout":
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFB)
        box.line.color.rgb = ACCENT_LT; box.line.width = Pt(1)
        box.text_frame.margin_left = Inches(0.12)
    tf = box.text_frame; tf.word_wrap = True
    for i, ln in enumerate(strip_inline(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = ""
        if kind == "bullet":
            prefix = ("    • " if level else "• ")
            # 번호불릿이면 자체 prefix 유지
            if re.match(r"^\d+\.", ln):
                prefix = ("    " if level else "")
        if prefix:
            r = p.add_run(); r.text = prefix; r.font.size = Pt(size); r.font.name = KR_FONT; r.font.color.rgb = color
        add_runs(p, ln, Pt(size), color=color, base_bold=bold)
    return top + Emu(int(h)) + Inches(0.12)


def main():
    with open(SRC, encoding="utf-8") as f:
        md = f.read()
    # 슬라이드 분할: '## 장표' ~ 다음 '## 장표'/'## 참고'
    parts = re.split(r"(?m)^## 장표 ", md)
    slide_chunks = parts[1:]  # [0]은 표지 전 인덱스
    # 마지막 청크에서 '## 참고 자료' 이후 잘라냄
    cleaned = []
    for ch in slide_chunks:
        ch = re.split(r"(?m)^## 참고 자료", ch)[0]
        cleaned.append(ch)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for chunk in cleaned:
        lines = chunk.split("\n")
        header = lines[0].strip()  # "N — Title"
        m = re.match(r"^(\d+)\s*—\s*(.*)$", header)
        if m:
            idx = m.group(1)
            title = m.group(2).strip()
        else:
            idx, title = "", header
        body_lines = lines[1:]
        blocks, notes = parse_blocks(body_lines)

        slide = prs.slides.add_slide(blank)
        add_title(slide, title, idx)
        y = CONTENT_TOP
        for b in blocks:
            if y > CONTENT_BOTTOM:
                pass  # 넘쳐도 계속 배치(사용자 후편집)
            if b[0] == "table":
                y = render_table(slide, b[1], y)
            elif b[0] == "code":
                y = render_code(slide, b[1], y)
            elif b[0] == "heading":
                y = render_text(slide, "heading", b[1], y)
            elif b[0] == "bullet":
                y = render_text(slide, "bullet", b[1], y, level=b[2])
            elif b[0] == "callout":
                y = render_text(slide, "callout", b[1], y)
            else:
                y = render_text(slide, "para", b[1], y)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(OUT)
    print(f"OK  slides={len(cleaned)}  ->  {OUT}")


if __name__ == "__main__":
    main()
