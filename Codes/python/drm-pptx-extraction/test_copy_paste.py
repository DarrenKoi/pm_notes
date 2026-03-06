"""Test script to verify if COM copy-paste from DRM PPTX to non-DRM PPTX works.

Run this first to check feasibility before using the full pipeline.

Usage:
    python test_copy_paste.py input/drm_file.pptx input/empty.pptx

Steps:
    1. Opens both files in PowerPoint via COM
    2. Tries to copy slide 1 from DRM → paste into non-DRM
    3. Saves the result and reports success/failure
    4. Attempts to read the saved file with python-pptx to verify it's not locked
"""

import argparse
import os
import shutil
import sys
import time
import traceback


def test_com_copy_paste(drm_path: str, non_drm_path: str):
    """Test if copy-paste from DRM to non-DRM works via PowerPoint COM."""
    drm_path = os.path.abspath(drm_path)
    non_drm_path = os.path.abspath(non_drm_path)
    output_path = os.path.abspath("output/test_copied.pptx")

    os.makedirs("output", exist_ok=True)

    print("=" * 60)
    print("Step 1: Check input files")
    print("=" * 60)
    for label, path in [("DRM PPTX", drm_path), ("Non-DRM PPTX", non_drm_path)]:
        if os.path.exists(path):
            size_kb = os.path.getsize(path) / 1024
            print(f"  [OK] {label}: {path} ({size_kb:.1f} KB)")
        else:
            print(f"  [FAIL] {label}: {path} NOT FOUND")
            sys.exit(1)

    # Copy template to output
    shutil.copy2(non_drm_path, output_path)

    print()
    print("=" * 60)
    print("Step 2: Open PowerPoint via COM and try copy-paste")
    print("=" * 60)

    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()
    ppt = None
    copy_success = False

    try:
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = True
        print("  [OK] PowerPoint launched")

        # Open DRM file
        print(f"  Opening DRM file (ReadOnly)...")
        try:
            drm_prs = ppt.Presentations.Open(drm_path, ReadOnly=True)
            print(f"  [OK] DRM file opened - {drm_prs.Slides.Count} slides")
        except Exception as e:
            print(f"  [FAIL] Cannot open DRM file: {e}")
            print("  -> DRM may block COM access entirely.")
            return

        # Open non-DRM output file
        print(f"  Opening output file...")
        try:
            out_prs = ppt.Presentations.Open(output_path)
            print(f"  [OK] Output file opened - {out_prs.Slides.Count} slides")
        except Exception as e:
            print(f"  [FAIL] Cannot open output file: {e}")
            drm_prs.Close()
            return

        # Delete existing slides in output
        while out_prs.Slides.Count > 0:
            out_prs.Slides(1).Delete()

        # Try copying just the first slide
        print()
        print("  Attempting to copy slide 1 from DRM...")
        try:
            drm_prs.Slides(1).Copy()
            time.sleep(0.5)
            print("  [OK] Slide copied to clipboard")
        except Exception as e:
            print(f"  [FAIL] Copy failed: {e}")
            print("  -> DRM may block clipboard copy via COM.")
            drm_prs.Close()
            out_prs.Close()
            return

        print("  Attempting to paste into output...")
        try:
            out_prs.Slides.Paste()
            time.sleep(0.5)
            print(f"  [OK] Paste succeeded - output now has {out_prs.Slides.Count} slides")
            copy_success = True
        except Exception as e:
            print(f"  [FAIL] Paste failed: {e}")
            print("  -> Paste may be blocked or clipboard was empty.")
            drm_prs.Close()
            out_prs.Close()
            return

        # Save
        print()
        print("  Saving output file...")
        try:
            out_prs.Save()
            print(f"  [OK] Saved to: {output_path}")
        except Exception as e:
            print(f"  [FAIL] Save failed: {e}")
            print("  -> The file may be locked or DRM may propagate to pasted content.")

        drm_prs.Close()
        out_prs.Close()

    except Exception as e:
        print(f"  [ERROR] Unexpected error: {e}")
        traceback.print_exc()
    finally:
        if ppt is not None:
            ppt.Quit()
        pythoncom.CoUninitialize()

    if not copy_success:
        return

    # Step 3: Verify the saved file is readable by python-pptx
    print()
    print("=" * 60)
    print("Step 3: Verify saved file is readable by python-pptx")
    print("=" * 60)
    test_extract(output_path)


def test_extract(pptx_path: str):
    """Try opening the file with python-pptx to check if it's locked or corrupted."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        print(f"  [OK] python-pptx opened successfully - {slide_count} slides")

        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            images = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text[:50])
                if hasattr(shape, "image"):
                    try:
                        _ = shape.image.blob
                        images += 1
                    except Exception:
                        pass
            print(f"  Slide {i}: {len(texts)} text blocks, {images} images")
            for t in texts[:3]:
                print(f"    - \"{t}\"")

        print()
        print("  [SUCCESS] File is fully readable. You can proceed with full extraction.")

    except Exception as e:
        print(f"  [FAIL] python-pptx cannot open file: {e}")
        print("  -> The file may be corrupted or still DRM-locked after copy-paste.")
        print("  -> Try opening the file manually in PowerPoint and re-saving as a new file.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Test DRM PPTX copy-paste feasibility")
    parser.add_argument("drm_pptx", help="Path to DRM-protected PPTX")
    parser.add_argument("non_drm_pptx", help="Path to empty non-DRM PPTX (template)")
    args = parser.parse_args()

    test_com_copy_paste(args.drm_pptx, args.non_drm_pptx)
