"""Extract text and images from a non-DRM PPTX file using python-pptx."""

import json
import os

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_pptx(pptx_path: str, output_dir: str) -> dict:
    """Extract text and images from each slide, save to output_dir.

    Args:
        pptx_path: Path to the (non-DRM) PPTX file.
        output_dir: Directory to write images and result JSON.

    Returns:
        Structured dict with slide-by-slide extraction results.
    """
    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)

    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    result = {
        "source": os.path.basename(pptx_path),
        "total_slides": len(prs.slides),
        "slides": [],
    }

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide_number": slide_idx,
            "title": None,
            "texts": [],
            "images": [],
        }

        for shape in slide.shapes:
            # Extract title
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                        slide_data["title"] = text
                    else:
                        slide_data["texts"].append(text)

            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                img_filename = f"slide_{slide_idx:02d}_img_{len(slide_data['images']) + 1:02d}.{ext}"
                img_path = os.path.join(images_dir, img_filename)

                with open(img_path, "wb") as f:
                    f.write(image.blob)

                # Store relative path from output_dir
                slide_data["images"].append(os.path.join("images", img_filename))
                print(f"  Saved image: {img_filename}")

            # Extract images from group shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child_shape in shape.shapes:
                    if child_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = child_shape.image
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        img_filename = f"slide_{slide_idx:02d}_img_{len(slide_data['images']) + 1:02d}.{ext}"
                        img_path = os.path.join(images_dir, img_filename)

                        with open(img_path, "wb") as f:
                            f.write(image.blob)

                        slide_data["images"].append(os.path.join("images", img_filename))
                        print(f"  Saved image: {img_filename}")

        # If title wasn't found via shape_id, use the first text as title
        if slide_data["title"] is None and slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_data["title"] = slide.shapes.title.text_frame.text.strip() or None

        print(f"[extract] Slide {slide_idx}: title={slide_data['title']!r}, "
              f"{len(slide_data['texts'])} text blocks, {len(slide_data['images'])} images")

        result["slides"].append(slide_data)

    # Save JSON result
    json_path = os.path.join(output_dir, "extraction_result.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"\n[extract] Result saved to: {json_path}")

    return result
